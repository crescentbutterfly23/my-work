# -*- coding: utf-8 -*-
"""Build the whole manual: SVG -> PNG -> PPTX -> PDF, in every language."""
import os, sys, glob, shutil, subprocess, urllib.parse
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

import core
from strings import LANGS
import pages1 as A
import pages2 as B

BUILD = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(BUILD)
TMP = BUILD + "/tmp"
EDGE = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"

PAGES = ([(i, getattr(A, "p%02d" % i)) for i in range(1, 19)] +
         [(i, getattr(B, "p%02d" % i)) for i in range(19, 37)])

NAMES = {1: "portada", 2: "contenido", 3: "sec-marca", 4: "origen", 5: "proposito", 6: "valores",
         7: "posicionamiento", 8: "personalidad", 9: "territorio", 10: "manifiesto", 11: "casas",
         12: "arquitectura", 13: "voz", 14: "sec-identidad", 15: "anatomia", 16: "construccion",
         17: "lockups", 18: "area-respeto", 19: "versiones-color", 20: "usos-incorrectos",
         21: "monograma", 22: "sec-color", 23: "paleta", 24: "proporcion", 25: "contraste",
         26: "sec-tipografia", 27: "familias", 28: "escala", 29: "sec-velo", 30: "velo",
         31: "texturas", 32: "imagen", 33: "sec-aplicaciones", 34: "papeleria", 35: "digital",
         36: "cierre"}

TITLES_ES = {1: "Portada", 2: "Contenido", 3: "01 · La marca", 4: "El origen", 5: "Propósito, visión y misión",
             6: "Valores", 7: "Posicionamiento", 8: "Personalidad", 9: "Territorio de marca",
             10: "Manifiesto", 11: "Las seis casas", 12: "Arquitectura de marca", 13: "Voz y tono",
             14: "02 · Identidad", 15: "Anatomía del logotipo", 16: "Construcción", 17: "Lockups",
             18: "Área de respeto", 19: "Versiones de color", 20: "Usos incorrectos", 21: "El monograma",
             22: "03 · Color", 23: "La paleta", 24: "Proporción", 25: "Contraste y accesibilidad",
             26: "04 · Tipografía", 27: "Las familias", 28: "Escala y jerarquía", 29: "05 · El velo",
             30: "El velo", 31: "Tramas y texturas", 32: "Dirección de imagen", 33: "06 · Aplicaciones",
             34: "Papelería", 35: "Digital", 36: "Cierre"}

TITLES_EN = {1: "Cover", 2: "Contents", 3: "01 · The brand", 4: "The origin", 5: "Purpose, vision, mission",
             6: "Values", 7: "Positioning", 8: "Personality", 9: "Brand territory", 10: "Manifesto",
             11: "The six houses", 12: "Brand architecture", 13: "Voice and tone", 14: "02 · Identity",
             15: "Anatomy of the logotype", 16: "Construction", 17: "Lockups", 18: "Clear space",
             19: "Colour versions", 20: "Misuse", 21: "The monogram", 22: "03 · Colour", 23: "The palette",
             24: "Proportion", 25: "Contrast and accessibility", 26: "04 · Typography", 27: "The families",
             28: "Scale and hierarchy", 29: "05 · The veil", 30: "The veil", 31: "Patterns and textures",
             32: "Image direction", 33: "06 · Applications", 34: "Stationery", 35: "Digital", 36: "Closing"}

def build_svgs(code):
    L = LANGS[code]
    out = ROOT + "/SVG_%s" % code
    os.makedirs(out, exist_ok=True)
    for f in glob.glob(out + "/*.svg"):
        os.remove(f)
    for n, fn in PAGES:
        p = out + "/velum-%s-%02d-%s.svg" % (code.lower(), n, NAMES[n])
        open(p, "w", encoding="utf-8").write(fn(L))
    print(code, "SVG:", len(os.listdir(out)))

def render(code):
    src = ROOT + "/SVG_%s" % code
    dst = ROOT + "/PNG_%s" % code
    os.makedirs(dst, exist_ok=True)
    os.makedirs(TMP, exist_ok=True)
    for f in glob.glob(dst + "/*.png"):
        os.remove(f)
    html = ('<!doctype html><meta charset="utf-8"><style>html,body{margin:0;padding:0;'
            'background:#F4F0EE;width:1920px;height:1080px;overflow:hidden}'
            'img{display:block;width:1920px;height:1080px}</style><img src="%s">')
    for f in sorted(glob.glob(src + "/*.svg")):
        base = os.path.splitext(os.path.basename(f))[0]
        shutil.copy(f, TMP + "/" + os.path.basename(f))
        h = TMP + "/r.html"
        open(h, "w", encoding="utf-8").write(html % urllib.parse.quote(os.path.basename(f)))
        out = dst + "/" + base + ".png"
        subprocess.run([EDGE, "--headless=new", "--disable-gpu", "--hide-scrollbars",
                        "--force-device-scale-factor=1", "--window-size=1920,1080",
                        "--screenshot=" + out, "file:///" + h.replace("\\", "/")],
                       capture_output=True, timeout=180)
    print(code, "PNG:", len(os.listdir(dst)))

def deck(code):
    from pptx import Presentation
    from pptx.util import Emu
    titles = TITLES_ES if code == "ES" else TITLES_EN
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]
    files = sorted(glob.glob(ROOT + "/PNG_%s/*.png" % code))
    for f in files:
        n = int(os.path.basename(f).split("-")[2])
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(f, 0, 0, width=prs.slide_width, height=prs.slide_height)
        s.notes_slide.notes_text_frame.text = "%02d — %s" % (n, titles.get(n, ""))
    name = "VELUM_Manual_de_Marca_ES.pptx" if code == "ES" else "VELUM_Brand_Guidelines_EN.pptx"
    prs.save(ROOT + "/" + name)
    print(code, "PPTX:", os.path.getsize(ROOT + "/" + name) // 1024, "KB")

def pdf(code):
    from PIL import Image
    files = sorted(glob.glob(ROOT + "/PNG_%s/*.png" % code))
    ims = [Image.open(f).convert("RGB") for f in files]
    name = "VELUM_Manual_de_Marca_ES.pdf" if code == "ES" else "VELUM_Brand_Guidelines_EN.pdf"
    ims[0].save(ROOT + "/" + name, save_all=True, append_images=ims[1:], resolution=144.0)
    print(code, "PDF:", os.path.getsize(ROOT + "/" + name) // 1024, "KB", len(ims), "pages")

def sheet(code):
    from PIL import Image
    files = sorted(glob.glob(ROOT + "/PNG_%s/*.png" % code))
    cols, rows = 6, 6
    s = Image.new("RGB", (330 * cols, 186 * rows), "#888")
    for k, f in enumerate(files):
        s.paste(Image.open(f).convert("RGB").resize((330, 186)), (330 * (k % cols), 186 * (k // cols)))
    s.save(TMP + "/sheet_%s.png" % code)

if __name__ == "__main__":
    codes = sys.argv[1:] or ["ES", "EN"]
    for c in codes:
        build_svgs(c)
        render(c)
        deck(c)
        pdf(c)
        sheet(c)
